
import base64
import io
import os
import tempfile
import logging
from typing import Dict, List, Any
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from pptx import Presentation
import cv2
import numpy as np

logger = logging.getLogger(__name__)

class SlideProcessor:
    def __init__(self):
        self.temp_dir = "temp"
        os.makedirs(self.temp_dir, exist_ok=True)
    
    async def extract_content(self, slide_data: str, filename: str) -> Dict[str, Any]:
        """Extract content from uploaded slides"""
        try:
            # Decode base64 data
            file_data = base64.b64decode(slide_data)
            
            # Determine file type
            file_ext = filename.lower().split('.')[-1]
            
            if file_ext in ['pptx', 'ppt']:
                return await self._extract_from_powerpoint(file_data)
            elif file_ext == 'pdf':
                return await self._extract_from_pdf(file_data)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
                
        except Exception as e:
            logger.error(f"Failed to extract slide content: {e}")
            raise
    
    async def _extract_from_powerpoint(self, file_data: bytes) -> Dict[str, Any]:
        """Extract content from PowerPoint files"""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            tmp_file.write(file_data)
            tmp_file.flush()
            
            try:
                presentation = Presentation(tmp_file.name)
                slides_content = []
                
                for i, slide in enumerate(presentation.slides):
                    slide_text = []
                    slide_images = []
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    slides_content.append({
                        'slide_number': i + 1,
                        'text': '\n'.join(slide_text),
                        'images': slide_images
                    })
                
                return {
                    'slides': slides_content,
                    'total_slides': len(slides_content),
                    'extracted_text': self._combine_slide_text(slides_content)
                }
                
            finally:
                os.unlink(tmp_file.name)
    
    async def _extract_from_pdf(self, file_data: bytes) -> Dict[str, Any]:
        """Extract content from PDF files"""
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
            tmp_file.write(file_data)
            tmp_file.flush()
            
            try:
                doc = fitz.open(tmp_file.name)
                slides_content = []
                
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    
                    # Extract text
                    text = page.get_text()
                    
                    # Extract images and convert to PIL format for OCR if needed
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    
                    # OCR on image if text extraction didn't work well
                    if len(text.strip()) < 50:  # If very little text found
                        image = Image.open(io.BytesIO(img_data))
                        ocr_text = pytesseract.image_to_string(image)
                        text = f"{text}\n{ocr_text}" if text.strip() else ocr_text
                    
                    slides_content.append({
                        'slide_number': page_num + 1,
                        'text': text.strip(),
                        'images': []
                    })
                
                doc.close()
                
                return {
                    'slides': slides_content,
                    'total_slides': len(slides_content),
                    'extracted_text': self._combine_slide_text(slides_content)
                }
                
            finally:
                os.unlink(tmp_file.name)
    
    def _combine_slide_text(self, slides_content: List[Dict]) -> str:
        """Combine all slide text into a coherent document"""
        combined_text = []
        for slide in slides_content:
            if slide['text'].strip():
                combined_text.append(f"Slide {slide['slide_number']}: {slide['text']}")
        
        return '\n\n'.join(combined_text)
    
    async def extract_video_content(self, video_base64: str) -> Dict[str, Any]:
        """Extract frames and audio from video for analysis"""
        try:
            # Decode video data
            video_data = base64.b64decode(video_base64)
            
            with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as tmp_file:
                tmp_file.write(video_data)
                tmp_file.flush()
                
                try:
                    # Extract frames using OpenCV
                    cap = cv2.VideoCapture(tmp_file.name)
                    frames = []
                    frame_count = 0
                    
                    # Extract every 30th frame for analysis
                    while True:
                        ret, frame = cap.read()
                        if not ret:
                            break
                        
                        if frame_count % 30 == 0:  # Every second at 30fps
                            # Convert to base64 for analysis
                            _, buffer = cv2.imencode('.jpg', frame)
                            frame_b64 = base64.b64encode(buffer).decode()
                            frames.append({
                                'timestamp': frame_count / 30.0,
                                'frame_data': frame_b64
                            })
                        
                        frame_count += 1
                    
                    cap.release()
                    
                    # Get video metadata
                    duration = frame_count / 30.0  # Assuming 30fps
                    
                    return {
                        'frames': frames[:10],  # Limit to first 10 frames for analysis
                        'duration': duration,
                        'frame_count': frame_count,
                        'fps': 30
                    }
                    
                finally:
                    os.unlink(tmp_file.name)
                    
        except Exception as e:
            logger.error(f"Failed to extract video content: {e}")
            raise
